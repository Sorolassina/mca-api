from pptx import Presentation
import re


def replace_ignore_case(text: str, replacements: dict) -> str:
    for old, new in replacements.items():
        pattern = re.compile(re.escape(old), re.IGNORECASE)
        text = pattern.sub(new, text)
    return text

def modifier_contenu_powerpoint(file_path: str, replacements: dict):
    try:
        prs = Presentation(file_path)
        modified = False

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        new_text = replace_ignore_case(run.text, replacements)
                        if new_text != run.text:
                            run.text = new_text
                            modified = True

        if modified:
            prs.save(file_path)

    except Exception as e:
        print(f"❌ Erreur traitement PowerPoint {file_path} : {e}")
